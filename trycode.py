import google.generativeai as genai
import os
import fitz  # PyMuPDF for PDF
import io
import pandas as pd  # for Excel and CSV
from docx import Document  # for .docx files
from pptx import Presentation  # for .ppt and .pptx files
import argparse
import sys
import logging
import json # Import json for potentially validating output later

# --- Set up logging ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# --- Text Extraction Function ---
def extract_text_from_file(file_path):
    logger.info(f"Attempting to extract text from file: {file_path}")
    if not os.path.exists(file_path):
        logger.error(f"File does not exist: {file_path}")
        return f"[Error: File {file_path} does not exist]" # Return an error message

    file_extension = os.path.splitext(file_path)[1].lower()
    text = ""

    try:
        with open(file_path, "rb") as file:
            file_bytes = file.read()

        if file_extension == ".pdf":
            logger.info("Processing PDF file")
            # Use memory stream
            pdf_document = fitz.open(stream=file_bytes, filetype="pdf")
            for page_num in range(pdf_document.page_count):
                page = pdf_document.load_page(page_num)
                text += page.get_text("text") + "\n"
            pdf_document.close()

        elif file_extension == ".docx":
            logger.info("Processing DOCX file")
            doc = Document(io.BytesIO(file_bytes))
            for para in doc.paragraphs:
                text += para.text + "\n"

        elif file_extension in [".ppt", ".pptx"]:
            logger.info("Processing PowerPoint file")
            prs = Presentation(io.BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        elif file_extension in [".xlsx", ".xls"]:
            logger.info("Processing Excel file")
            # Specify engine='openpyxl' or 'xlrd' if needed based on pandas version/installation
            excel_data = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
            for sheet_name, df in excel_data.items():
                text += f"--- Sheet: {sheet_name} ---\n"
                # Convert DataFrame to a more readable text format, maybe not full CSV for prompt
                # Using tab separated or simple string representation
                text += df.to_string(index=False) + "\n\n" # Changed from to_csv
                # Or if CSV is needed for structure:
                # text += df.to_csv(index=False) + "\n\n"


        elif file_extension in [".txt", ".csv"]:
            logger.info("Processing TXT or CSV file")
            text = file_bytes.decode('utf-8', errors='ignore')

        else:
            logger.warning(f"Unsupported file type: {file_extension}. Attempting to read as text")
            try:
                text = file_bytes.decode('utf-8', errors='ignore')
            except Exception as e:
                logger.error(f"Failed to process {file_path} as text: {e}")
                return f"[Error: Unsupported file type {file_extension}. Could not extract text: {e}]"

    except Exception as e:
        logger.error(f"Error processing {file_path}: {e}")
        return f"[Error processing {file_path}: {e}]" # Return an error message

    if not text.strip(): # Also check for errors returned by the function itself
         if not text.startswith("[Error") and not text.startswith("[Warning"):
            logger.warning(f"No readable text extracted from {file_path}. File might be scanned, empty, or have complex formatting.")
            return f"[Warning: No readable text extracted from {file_path}. File might be scanned, empty, or have complex formatting.]"

    logger.info(f"Successfully extracted text from {file_path}")
    return text

# --- Base Prompt Instructions (includes placeholders for user data) ---
BASE_PROMPT_INSTRUCTIONS = """
Role: Bewell AI Assistant.
Persona: Holistic women's health expert, precision medicine, functional doctor, Women focused care.
Tone: Approachable, professional, empathetic, supportive, clear, accessible. Avoid casual language. Do not use the pronoun "I". Avoid personifying the analysis tone and manner. Focus on empowering the user with clear, accurate, and personalized insights from Bewell.

Input: User's health assessment text and lab report text. Analyze *only* the information provided in these texts. Base the analysis, recommendations, and rationale *solely* on the specific biomarkers and symptoms reported by the user in the provided data.

---

User Data:

Here is the user's Lab Report text:
{lab_report_text}

Here is the user's Health Assessment text:
{health_assessment_text}

---

Instructions for Output:

Generate ONE complete response as a JSON object. Do NOT include any introductory or concluding text outside the JSON object. Populate all fields based *only* on the provided user data, adhering to the persona, tone, and data constraints specified. Calculate the biomarker counts based on the provided lab results and ranges. Compile the lists in the 'action_plan' by extracting relevant information from the detailed sections as described in the JSON structure definition below. Ensure all string values within the JSON are plain text; do NOT use Markdown formatting (like **, #, -) within the JSON values.

Here is the required JSON object structure:
"""

# --- JSON Structure Definition (this string *is* the structure definition for the AI) ---
JSON_STRUCTURE_DEFINITION = """
{
  "lab_analysis": {
    "overall_summary": "string - Synthesize user's current health status based solely on provided lab results and health assessment. Highlight significant areas of concern or strengths in simple, personalized language.",
    "biomarkers_tested_count": "integer - Total number of biomarkers listed in the provided lab report.",
    "biomarker_categories_summary": {
      "description": "Count of biomarkers categorized by status based on result vs. range.",
      "structure": {
        "optimal": "integer - Count of biomarkers whose status is 'Optimal (Green)'.",
        "keep_in_mind": "integer - Count of biomarkers whose status is 'Keep in Mind (Yellow)'.",
        "attention_needed": "integer - Count of biomarkers whose status is 'Attention Needed (Orange-Red)'."
      }
    },
    "detailed_biomarkers": "array of objects - Contains details for *each* biomarker listed in the provided lab report.",
    "detailed_biomarkers_item_structure": {
      "name": "string - The full name of the biomarker.",
      "status": "string - Machine-readable status: 'optimal', 'keep_in_mind', or 'attention_needed'. Determine this based on the result relative to the provided range.",
      "status_label": "string - Display label mirroring original status labels (e.g., 'Optimal (Green)', 'Keep in Mind (Yellow)', 'Attention Needed (Orange-Red)').",
      "result": "string - User's specific measured result, including units.",
      "range": "string - The provided reference or optimal range.",
      "cycle_impact": "string - Detail any known fluctuations or impacts specific menstrual cycle phases have, or state 'Not typically impacted by cycle' / 'Cycle impact not well-established'.",
      "why_it_matters": "string - Explanation of the biomarker's primary function, its importance specifically to women's health, and the potential practical implications if *this user's specific level* is abnormal or borderline. Use clear, science-backed explanations without medical jargon."
    },
    "crucial_biomarkers_to_measure": "array of objects - A list of essential biomarkers women should measure regularly.",
     "crucial_biomarkers_item_structure": {
        "name": "string - The name of the crucial biomarker.",
        "importance": "string - Brief explanation of its importance in accessible language."
     },
    "health_recommendation_summary": "array of strings - Clear, concise, *actionable* steps tailored *specifically* to the user's *provided* lab results and health assessment findings."
  },
  "four_pillars": {
    "introduction": "string - Briefly summarize the user's overall health based on findings from the lab analysis section in clear, accessible language.",
    "pillars": "array of objects - Contains the detailed analysis for each of the four pillars: Eat Well, Sleep Well, Move Well, Recover Well.",
    "pillars_item_structure": {
      "name": "string - The name of the pillar (e.g., 'Eat Well', 'Sleep Well', 'Move Well', 'Recover Well').",
      "why_it_matters": "string - Explain how this pillar is *specifically relevant to this user's unique health assessment details and lab findings*. Use reachable language.",
      "personalized_recommendations": "array of strings - Provide *actionable, personalized* advice tailored to *this user's specific status and lifestyle*. Recommendations must be achievable.",
      "root_cause_correlation": "string - Clearly explain in accessible language how *each recommendation* connects directly to the *root causes or contributing factors* identified *in this user's lab results and health assessment*.",
      "science_based_explanation": "string - For *each recommendation*, provide a clear, simple scientific basis focused on practical user benefits, without medical jargon.",
      "additional_guidance": {
         "description": "Specific lists of recommended/avoided items for this pillar.",
         "structure": {
            "recommended_foods": "array of objects - For 'Eat Well': lists of recommended top foods. Each object: {name: string, description: string|null}.",
            "cautious_foods": "array of objects - For 'Eat Well': lists of foods to approach cautiously. Each object: {name: string, description: string|null}.",
            "recommended_workouts": "array of objects - For 'Move Well': list of top recommended workouts. Each object: {name: string, description: string|null}.",
            "avoid_habits_move": "array of objects - For 'Move Well': list of habits to avoid. Each object: {name: string, description: string|null}.",
            "recommended_recovery_tips": "array of objects - For 'Sleep Well' and 'Recover Well': list of top recommended recovery tips. Each object: {name: string, description: string|null}.",
            "avoid_habits_rest_recover": "array of objects - For 'Sleep Well' and 'Recover Well': list of habits to avoid. Each object: {name: string, description: string|null}."
         }
      }
    }
  },
  "supplements": {
    "description": "Personalized supplement recommendations based on provided data.",
    "structure": {
      "recommendations": "array of objects - Details for each recommended supplement.",
      "recommendations_item_structure": {
        "name": "string - Supplement Name.",
        "rationale": "string - Personalized Rationale: Clearly explain *why* recommended based on user's biomarkers and reported health assessment symptoms. Explain how it addresses *this user's specific issues reported*. Use simple, accessible language.",
        "expected_outcomes": "string - Expected Outcomes: Describe tangible, *personalized* benefits the user can realistically notice.",
        "dosage_and_timing": "string - Recommended Dosage & Timing: Clearly outline precise dosage instructions and optimal timing.",
        "situational_cyclical_considerations": "string - Situational/Cyclical Considerations: Clearly identify if beneficial during specific menstrual cycle phases or particular life circumstances *if applicable and relevant to the supplement and user's provided profile*. Explain *why* this is the case simply."
      },
      "conclusion": "string - The concluding guidance: Concise, reassuring guidance and clear statement about supplements being adjunctive and requiring medical consultation."
    }
  },
   "action_plan": {
      "description": "Consolidated actionable recommendations for quick reference on the main summary.",
      "structure": {
          "foods_to_enjoy": "array of strings - Extract and consolidate the 'name' and optionally 'description' (or just 'name') from the 'recommended_foods' list under the 'Eat Well' pillar.",
          "foods_to_limit": "array of strings - Extract and consolidate the 'name' and optionally 'description' (or just 'name') from the 'cautious_foods' list under the 'Eat Well' pillar.",
          "daily_habits": "array of strings - Extract and consolidate key actionable daily practices mentioned in the 'Health Recommendation Summary' and 'personalized_recommendations' across the Four Pillars (e.g., related to consistent sleep, hydration, meal timing).",
          "rest_and_recovery": "array of strings - Extract and consolidate the 'name' and optionally 'description' (or just 'name') from the 'recommended_recovery_tips' lists under the 'Sleep Well' and 'Recover Well' pillars.",
          "movement": "array of strings - Extract and consolidate the 'name' and optionally 'description' (or just 'name') from the 'recommended_workouts' list under the 'Move Well' pillar."
      }
   }
}
"""


def main():
    # --- Command-Line Argument Parsing ---
    parser = argparse.ArgumentParser(description="Bewell AI Health Analyzer: Generate health analysis from health assessment (required) and lab report (optional) files.")
    parser.add_argument("--lab-report", type=str, help="Path to lab report file (any format) [optional]")
    parser.add_argument("--health-assessment", type=str, required=True, help="Path to health assessment file (any format) [required]")
    args = parser.parse_args()

    logger.info("Starting health analysis process")
    logger.info(f"Health assessment file: {args.health_assessment}")
    logger.info(f"Lab report file: {args.lab_report if args.lab_report else 'Not provided'}")

    # --- Extract Text from Files ---
    # Check if files exist before trying to extract
    if args.lab_report and not os.path.exists(args.lab_report):
         logger.error(f"Lab report file not found: {args.lab_report}")
         print(f"Error: Lab report file not found at {args.lab_report}")
         sys.exit(1)

    if not os.path.exists(args.health_assessment):
         logger.error(f"Health assessment file not found: {args.health_assessment}")
         print(f"Error: Health assessment file not found at {args.health_assessment}")
         sys.exit(1)


    lab_report_text = extract_text_from_file(args.lab_report) if args.lab_report else ""
    health_assessment_text = extract_text_from_file(args.health_assessment)

    # --- Validate Extracted Text ---
    # Check if extraction resulted in an error message
    if health_assessment_text.startswith("[Error"):
        logger.error(f"Failed to extract valid text from health assessment file: {health_assessment_text}")
        print(f"Error: Failed to extract valid text from health assessment file: {health_assessment_text}")
        sys.exit(1)
    if args.lab_report and lab_report_text.startswith("[Error"):
        logger.error(f"Failed to extract valid text from lab report file: {lab_report_text}")
        print(f"Error: Failed to extract valid text from lab report file: {lab_report_text}")
        sys.exit(1)


    if not health_assessment_text.strip() or health_assessment_text.startswith("[Warning"):
        logger.warning(f"Health assessment text seems empty or problematic: {health_assessment_text}")
        print(f"Warning: Health assessment text seems empty or problematic. Analysis quality may be affected.")
        # Don't exit, allow analysis with potentially limited data

    if args.lab_report and (not lab_report_text.strip() or lab_report_text.startswith("[Warning")):
        logger.warning(f"Lab report text seems empty or problematic: {lab_report_text}")
        print(f"Warning: Lab report text seems empty or problematic. Analysis quality may be affected.")
        # Don't exit, allow analysis with potentially limited data


    # --- API Key and Model Configuration ---
    # Use environment variable first, then hardcoded as fallback
    api_key = "AIzaSyArQ9zeya1SO-IwsMappkLStXYT0W7WXfk" # Replace with your fallback key or remove fallback

    if not api_key:
        logger.error("Google/Gemini API key not found")
        print("Error: Google/Gemini API key not found. Set it in the script or as an environment variable (GOOGLE_API_KEY).")
        sys.exit(1)

    model_name = "gemini-2.0-flash" # Or "gemini-1.5-flash-latest" or "gemini-1.5-pro-latest"
    generation_config = {
        "temperature": 0.7,
        "top_p": 0.95,
        "top_k": 64,
        "max_output_tokens": 8192,
        # response_mime_type should ideally be 'application/json' if the model supports it well.
        # For older models or inconsistent JSON output, keeping it 'text/plain'
        # and parsing the string response in code is safer.
        # Let's stick to text/plain and parse for robustness.
        "response_mime_type": "text/plain",
    }

    # --- Configure Gemini API ---
    try:
        logger.info("Configuring Gemini API")
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel(
            model_name=model_name,
            generation_config=generation_config,
        )

        # --- Format Prompt ---
        logger.info("Formatting prompt with user data and JSON structure instructions")
        # Format the base instructions string with the actual user data
        formatted_instructions = BASE_PROMPT_INSTRUCTIONS.format(
            lab_report_text=lab_report_text,
            health_assessment_text=health_assessment_text
        )
        # Combine the formatted instructions with the JSON structure definition
        prompt = formatted_instructions + JSON_STRUCTURE_DEFINITION

        # Optional: Print the full prompt for debugging
        # print("\n--- FULL PROMPT SENT TO AI ---\n")
        # print(prompt)
        # print("\n-----------------------------\n")


        print("Generating health analysis... This may take a few moments.")
        logger.info("Generating health analysis")

        # --- Generate Analysis ---
        response = model.generate_content(prompt)

        # --- Output Results ---
        print("\n=== Your Personalized Bewell Analysis (JSON Output) ===\n")
        if response and response.text:
            try:
                # Attempt to parse the response text as JSON
                analysis_data = json.loads(response.text)
                # Print pretty-printed JSON
                print(json.dumps(analysis_data, indent=2))
                logger.info("Analysis generated and parsed as JSON successfully")

            except json.JSONDecodeError as e:
                logger.error(f"Failed to parse AI response as JSON: {e}")
                print(f"Error: Failed to parse AI response as JSON. The model may not have followed the structure instructions.")
                print("Raw response text:")
                print(response.text) # Print raw text if JSON parsing fails
                sys.exit(1)

        elif response and response.prompt_feedback:
            logger.error(f"Analysis request blocked. Reason: {response.prompt_feedback.block_reason}")
            print(f"Error: Analysis request was blocked by the model. Reason: {response.prompt_feedback.block_reason}")
            if response.prompt_feedback.safety_ratings:
                print(f"Safety Ratings: {response.prompt_feedback.safety_ratings}")
            print("Please review your input files for sensitive content and try again.")
            sys.exit(1)
        else:
            logger.error("Failed to generate a valid response from the model (empty response)")
            print("Error: Failed to generate a valid response from the model (response was empty).")
            sys.exit(1)

        # --- Output Raw Extracted Text (Optional - keep for debugging) ---
        print("\n=== Raw Input Sent to AI (for debugging) ===\n")
        if lab_report_text.strip() and not lab_report_text.startswith("[Error"):
            print("Lab Report Text:")
            print(lab_report_text)
            print("\n")
        else:
            print("No valid Lab Report text provided or extracted.")

        if health_assessment_text.strip() and not health_assessment_text.startswith("[Error"):
            print("Health Assessment Text:")
            print(health_assessment_text)
            print("\n")
        else:
             print("No valid Health Assessment text provided or extracted.")

        logger.info("Raw input text output completed")

    except Exception as e:
        logger.error(f"An unexpected error occurred: {e}", exc_info=True) # Log traceback
        print(f"An unexpected error occurred: {e}")
        sys.exit(1)

    print("\n=== Footer ===")
    print("This analysis is AI-generated for informational purposes only. Always consult a medical professional for diagnosis and treatment.")
    logger.info("Health analysis process completed")

if __name__ == "__main__":
    main()